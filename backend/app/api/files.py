"""
Files API — Upload, extract text, list, delete
"""
from fastapi import APIRouter, UploadFile, File, HTTPException
import os, shutil, json, time
from pathlib import Path

router = APIRouter()

UPLOAD_DIR = os.path.join(os.path.dirname(__file__), "../../uploads")
META_FILE  = os.path.join(UPLOAD_DIR, "_meta.json")
os.makedirs(UPLOAD_DIR, exist_ok=True)

ALLOWED = {".pdf",".docx",".doc",".pptx",".ppt",".txt",".md",".xlsx",".xls",".png",".jpg",".jpeg"}


def _load():
    if os.path.exists(META_FILE):
        try:
            with open(META_FILE) as f: return json.load(f)
        except: pass
    return {}

def _save(meta):
    with open(META_FILE, "w") as f: json.dump(meta, f, indent=2)

def _fmt(b):
    if b < 1024: return f"{b}B"
    if b < 1048576: return f"{b//1024}KB"
    return f"{b/1048576:.1f}MB"

def _extract(path, filename):
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".pdf":
            try:
                import pdfplumber
                pages = []
                with pdfplumber.open(path) as pdf:
                    for page in pdf.pages[:25]:
                        t = page.extract_text()
                        if t: pages.append(t)
                text = "\n\n".join(pages)
                if text.strip(): return text[:8000]
            except: pass
            try:
                from pypdf import PdfReader
                r = PdfReader(path)
                pages = [p.extract_text() for p in r.pages[:25] if p.extract_text()]
                return "\n\n".join(pages)[:8000]
            except Exception as e:
                return f"[PDF extraction failed: {e}]"

        elif ext in (".docx", ".doc"):
            from docx import Document
            doc  = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text[:8000]

        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs    = Presentation(path)
            slides = []
            for i, slide in enumerate(prs.slides):
                texts = [s.text.strip() for s in slide.shapes if hasattr(s,"text") and s.text.strip()]
                if texts: slides.append(f"[Slide {i+1}] " + " | ".join(texts))
            return "\n".join(slides)[:8000]

        elif ext in (".xlsx", ".xls"):
            import openpyxl
            wb   = openpyxl.load_workbook(path, read_only=True, data_only=True)
            rows = []
            for name in wb.sheetnames[:3]:
                ws = wb[name]
                rows.append(f"[Sheet: {name}]")
                for row in ws.iter_rows(max_row=80, values_only=True):
                    r = " | ".join(str(c) for c in row if c is not None)
                    if r.strip(): rows.append(r)
            return "\n".join(rows)[:8000]

        elif ext in (".png", ".jpg", ".jpeg"):
            return "[Image file — OCR not available on free tier]"

        else:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(8000)

    except Exception as e:
        return f"[Extraction failed: {e}]"


@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED:
        raise HTTPException(400, f"File type '{ext}' not supported")

    safe = file.filename.replace(" ", "_")
    dest = os.path.join(UPLOAD_DIR, safe)
    with open(dest, "wb") as f:
        shutil.copyfileobj(file.file, f)

    text = _extract(dest, file.filename)
    size = os.path.getsize(dest)
    fid  = f"f_{abs(hash(safe+str(size)))}"

    meta = _load()
    meta[fid] = {
        "id": fid, "name": file.filename, "path": dest,
        "size": size, "size_str": _fmt(size),
        "ext": ext.lstrip("."), "text": text, "added": time.time(),
    }
    _save(meta)

    return {
        "id": fid, "name": file.filename,
        "size": _fmt(size), "ext": ext.lstrip("."),
        "chars": len(text),
        "preview": text[:200] + "..." if len(text) > 200 else text,
    }


@router.get("/list")
async def list_files():
    meta  = _load()
    files = [
        {"id": v["id"], "name": v["name"], "size": v.get("size_str","—"),
         "ext": v.get("ext",""), "chars": len(v.get("text","")), "added": v.get("added",0)}
        for k, v in meta.items() if not k.startswith("_")
    ]
    return {"files": sorted(files, key=lambda x: x["added"], reverse=True)}


@router.get("/{file_id}/text")
async def get_text(file_id: str):
    meta = _load()
    if file_id not in meta:
        raise HTTPException(404, "File not found")
    return {"text": meta[file_id]["text"], "name": meta[file_id]["name"]}


@router.delete("/{file_id}")
async def delete_file(file_id: str):
    meta = _load()
    if file_id not in meta:
        raise HTTPException(404, "File not found")
    try: os.remove(meta[file_id]["path"])
    except: pass
    del meta[file_id]
    _save(meta)
    return {"deleted": file_id}
