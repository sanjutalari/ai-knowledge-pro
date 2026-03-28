"""
Document Processing Service
Extracts text from PDF, DOCX, PPTX, TXT, Images (OCR)
All free/open-source libraries
"""
import os
import io
import re
from pathlib import Path
from typing import Optional


def extract_text(file_path: str, filename: str) -> str:
    """
    Main dispatcher — extracts text from any supported file type.
    Returns extracted text string.
    """
    ext = Path(filename).suffix.lower()

    extractors = {
        ".pdf":  _extract_pdf,
        ".docx": _extract_docx,
        ".doc":  _extract_docx,
        ".pptx": _extract_pptx,
        ".ppt":  _extract_pptx,
        ".txt":  _extract_txt,
        ".md":   _extract_txt,
        ".xlsx": _extract_xlsx,
        ".xls":  _extract_xlsx,
        ".png":  _extract_image_ocr,
        ".jpg":  _extract_image_ocr,
        ".jpeg": _extract_image_ocr,
        ".webp": _extract_image_ocr,
    }

    extractor = extractors.get(ext, _extract_txt)
    try:
        text = extractor(file_path)
        return _clean_text(text)
    except Exception as e:
        return f"[Could not extract text from {filename}: {e}]"


def _extract_pdf(path: str) -> str:
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages[:30]:  # limit to 30 pages
                text = page.extract_text()
                if text:
                    pages.append(text)
        if pages:
            return "\n\n".join(pages)
    except ImportError:
        pass

    # Fallback: pypdf
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        pages = []
        for page in reader.pages[:30]:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)
    except Exception as e:
        return f"[PDF extraction failed: {e}]"


def _extract_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        return "\n".join(paragraphs)
    except Exception as e:
        return f"[DOCX extraction failed: {e}]"


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[PPTX extraction failed: {e}]"


def _extract_xlsx(path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        rows = []
        for sheet in wb.sheetnames[:3]:  # limit to 3 sheets
            ws = wb[sheet]
            rows.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(max_row=100, values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    rows.append(row_text)
        return "\n".join(rows)
    except Exception as e:
        return f"[XLSX extraction failed: {e}]"


def _extract_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(50000)  # limit 50k chars
    except Exception as e:
        return f"[Text extraction failed: {e}]"


def _extract_image_ocr(path: str) -> str:
    try:
        import pytesseract
        from PIL import Image
        img  = Image.open(path)
        text = pytesseract.image_to_string(img)
        return text if text.strip() else "[No text found in image]"
    except ImportError:
        return "[OCR not available: install pytesseract and Pillow]"
    except Exception as e:
        return f"[OCR failed: {e}]"


def _clean_text(text: str) -> str:
    """Clean and truncate extracted text."""
    # Remove excessive whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {3,}', ' ', text)
    text = text.strip()
    # Limit to ~8000 chars to keep API calls fast & cheap
    if len(text) > 8000:
        text = text[:8000] + "\n\n[... document truncated for processing ...]"
    return text


def get_file_info(path: str, filename: str) -> dict:
    """Return basic file metadata."""
    size = os.path.getsize(path) if os.path.exists(path) else 0
    ext  = Path(filename).suffix.lower().lstrip(".")
    return {
        "name":      filename,
        "size":      size,
        "extension": ext,
        "size_str":  _fmt_size(size),
    }


def _fmt_size(b: int) -> str:
    if b < 1024:       return f"{b}B"
    if b < 1048576:    return f"{b//1024}KB"
    return f"{b/1048576:.1f}MB"
